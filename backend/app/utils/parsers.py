"""Document parsers for the ingestion pipeline (Phase 1 MVP)."""

from __future__ import annotations

import io
import os
import tempfile
from dataclasses import dataclass, field
from typing import Any

from app.core.config import settings

AUDIO_EXTENSIONS = frozenset({".mp3", ".mp4", ".m4a", ".wav", ".webm", ".ogg"})
AUDIO_MIME_TYPES = frozenset(
    {
        "audio/mpeg",
        "audio/mp3",
        "audio/mp4",
        "audio/x-m4a",
        "audio/wav",
        "audio/x-wav",
        "audio/webm",
        "audio/ogg",
        "video/mp4",
    }
)


@dataclass
class ParsedDocument:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    parser: str = "unknown"


def _word_count(text: str) -> int:
    return len(text.split()) if text.strip() else 0


def _docx_to_text(doc: Any) -> str:
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            prefix = ""
            if para.style.name.startswith("Heading"):
                level = para.style.name.split(" ")[-1]
                prefix = f"{'#' * int(level)} " if level.isdigit() else "## "
            parts.append(f"{prefix}{para.text.strip()}")

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


def parse_docx(content: bytes) -> str:
    """Парсить .docx через python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    return _docx_to_text(doc)


def _pptx_to_text(prs: Any) -> str:
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {slide_num}"]

        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_parts.append(f"### {slide.shapes.title.text.strip()}")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:
                    slide_parts.append(shape.text.strip())

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Notes: {notes_text}]")

        parts.append("\n".join(slide_parts))

    return "\n\n---\n\n".join(parts)


def parse_pptx(content: bytes) -> str:
    """Парсить .pptx через python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    return _pptx_to_text(prs)


def parse_audio(content: bytes, filename: str) -> str:
    """Транскрибує аудіо через OpenAI Whisper API."""
    from openai import OpenAI

    if not settings.openai_configured:
        raise ValueError("OpenAI API key is not configured for audio transcription")

    client = OpenAI(api_key=settings.OPENAI_API_KEY)

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "mp3"
    supported = {"mp3", "mp4", "mpeg", "mpga", "m4a", "wav", "webm", "ogg"}
    if ext not in supported:
        raise ValueError(f"Unsupported audio format: {ext}")

    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        with open(tmp_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file,
                response_format="text",
            )
        return f"[Audio Transcript]\n\n{transcript}"
    finally:
        os.unlink(tmp_path)


def parse_document(content: bytes, mime_type: str, filename: str = "") -> ParsedDocument:
    """
    Parse raw file bytes into plain text.

    Supports: PDF (PyMuPDF), TXT, Markdown, Excel (pandas), DOCX, PPTX, audio (Whisper).
    """
    mime = (mime_type or "").lower()
    name = (filename or "").lower()
    ext = ""
    if "." in name:
        ext = name[name.rfind(".") :]

    if mime == "application/pdf" or name.endswith(".pdf"):
        return _parse_pdf(content)
    if mime in ("text/plain",) or name.endswith(".txt"):
        return _parse_text(content, parser="plaintext")
    if mime in ("text/markdown",) or name.endswith((".md", ".markdown")):
        return _parse_text(content, parser="markdown")
    if (
        mime
        in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        )
        or name.endswith((".xlsx", ".xls"))
    ):
        return _parse_excel(content, filename)
    if (
        mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or name.endswith(".docx")
    ):
        return _parse_docx(content)
    if (
        mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or name.endswith(".pptx")
    ):
        return _parse_pptx(content)

    if ext in AUDIO_EXTENSIONS or mime in AUDIO_MIME_TYPES:
        return _parse_audio(content, filename)

    try:
        return _parse_text(content, parser="utf8_fallback")
    except Exception as exc:
        raise ValueError(f"Unsupported document type: {mime_type or filename}") from exc


def _parse_pdf(content: bytes) -> ParsedDocument:
    import fitz  # PyMuPDF

    pages: list[str] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page in doc:
            pages.append(page.get_text("text"))

    text = "\n\n".join(p.strip() for p in pages if p.strip())
    return ParsedDocument(
        text=text,
        metadata={
            "format": "pdf",
            "page_count": len(pages),
            "word_count": _word_count(text),
        },
        parser="pymupdf",
    )


def _parse_text(content: bytes, parser: str = "plaintext") -> ParsedDocument:
    text = content.decode("utf-8", errors="replace").strip()
    return ParsedDocument(
        text=text,
        metadata={"format": parser, "word_count": _word_count(text)},
        parser=parser,
    )


def _parse_excel(content: bytes, filename: str) -> ParsedDocument:
    import pandas as pd

    buffer = io.BytesIO(content)
    sheets = pd.read_excel(buffer, sheet_name=None, engine="openpyxl")

    parts: list[str] = []
    sheet_names: list[str] = []
    for sheet_name, df in sheets.items():
        sheet_names.append(str(sheet_name))
        parts.append(f"## Sheet: {sheet_name}\n{df.to_string(index=False)}")

    text = "\n\n".join(parts).strip()
    return ParsedDocument(
        text=text,
        metadata={
            "format": "excel",
            "sheet_count": len(sheets),
            "sheets": sheet_names,
            "word_count": _word_count(text),
        },
        parser="pandas",
    )


def _parse_docx(content: bytes) -> ParsedDocument:
    from docx import Document

    doc = Document(io.BytesIO(content))
    has_tables = len(doc.tables) > 0
    text = _docx_to_text(doc)

    return ParsedDocument(
        text=text,
        metadata={
            "format": "docx",
            "word_count": _word_count(text),
            "has_tables": has_tables,
        },
        parser="python-docx",
    )


def _parse_pptx(content: bytes) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slide_count = len(prs.slides)
    text = _pptx_to_text(prs)

    return ParsedDocument(
        text=text,
        metadata={
            "format": "pptx",
            "slide_count": slide_count,
            "word_count": _word_count(text),
        },
        parser="python-pptx",
    )


def _parse_audio(content: bytes, filename: str) -> ParsedDocument:
    text = parse_audio(content, filename)
    return ParsedDocument(
        text=text,
        metadata={
            "format": "audio",
            "word_count": _word_count(text),
        },
        parser="whisper-1",
    )
