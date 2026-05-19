"""Unit tests for extended document parsers (DOCX, PPTX, audio)."""

from __future__ import annotations

import io

import pytest
from docx import Document
from pptx import Presentation

from app.core.config import settings
from app.services.document_upload_policy import ALLOWED_EXTENSIONS
from app.utils.parsers import parse_audio, parse_docx, parse_document, parse_pptx


def _docx_bytes(build) -> bytes:
    doc = Document()
    build(doc)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(build) -> bytes:
    prs = Presentation()
    build(prs)
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def test_parse_docx_basic():
    def build(doc: Document) -> None:
        doc.add_paragraph("Hello World")
        doc.add_paragraph("Title", style="Heading 1")

    text = parse_docx(_docx_bytes(build))

    assert "Hello World" in text
    assert "# Title" in text


def test_parse_pptx_basic():
    def build(prs: Presentation) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Test Slide"
        slide.placeholders[1].text = "Content here"

    text = parse_pptx(_pptx_bytes(build))

    assert "## Slide 1" in text
    assert "Test Slide" in text
    assert "Content here" in text


def test_parse_docx_with_table():
    def build(doc: Document) -> None:
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"

    text = parse_docx(_docx_bytes(build))

    assert "A1 | B1" in text
    assert "A2 | B2" in text


def test_parse_pptx_multiple_slides():
    def build(prs: Presentation) -> None:
        layout = prs.slide_layouts[1]
        for i in range(1, 4):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i} title"

    text = parse_pptx(_pptx_bytes(build))

    assert "## Slide 1" in text
    assert "## Slide 2" in text
    assert "## Slide 3" in text
    assert "\n\n---\n\n" in text


def test_parse_audio_disabled(monkeypatch):
    monkeypatch.setattr(settings, "OPENAI_API_KEY", "")
    monkeypatch.setattr(settings, "AUDIO_TRANSCRIPTION_ENABLED", False)

    with pytest.raises(ValueError, match="OpenAI API key"):
        parse_audio(b"fake-audio", "sample.mp3")

    with pytest.raises(ValueError, match="OpenAI API key"):
        parse_document(b"fake-audio", "audio/mpeg", "sample.mp3")


def test_allowed_extensions_docx():
    assert ".docx" in ALLOWED_EXTENSIONS


def test_allowed_extensions_pptx():
    assert ".pptx" in ALLOWED_EXTENSIONS
